from pathlib import Path

from fastapi import HTTPException


SUPPORTED_OFFICE_EXTENSIONS = {".docx", ".pptx"}
LEGACY_OFFICE_EXTENSIONS = {".doc", ".ppt"}


def extract_docx_text(path: Path) -> tuple[str, int]:
    try:
        from docx import Document
    except ImportError as exc:
        raise HTTPException(
            status_code=500,
            detail="DOCX parsing dependency is not installed. Install python-docx.",
        ) from exc

    try:
        document = Document(str(path))
    except Exception as exc:
        raise HTTPException(status_code=400, detail="Corrupted or unreadable DOCX document.") from exc

    parts: list[str] = []
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if text:
            parts.append(text)

    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))

    return "\n\n".join(parts), len(document.paragraphs)


def _shape_text(shape) -> list[str]:
    texts: list[str] = []
    if getattr(shape, "has_text_frame", False):
        text = "\n".join(
            paragraph.text.strip()
            for paragraph in shape.text_frame.paragraphs
            if paragraph.text.strip()
        ).strip()
        if text:
            texts.append(text)

    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                texts.append(" | ".join(cells))

    for child in getattr(shape, "shapes", []):
        texts.extend(_shape_text(child))

    return texts


def extract_pptx_text(path: Path) -> tuple[str, int]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise HTTPException(
            status_code=500,
            detail="PPTX parsing dependency is not installed. Install python-pptx.",
        ) from exc

    try:
        presentation = Presentation(str(path))
    except Exception as exc:
        raise HTTPException(status_code=400, detail="Corrupted or unreadable PPTX presentation.") from exc

    slide_texts: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            texts.extend(_shape_text(shape))
        if texts:
            slide_texts.append(f"Slide {index}\n" + "\n".join(texts))

    return "\n\n".join(slide_texts), len(presentation.slides)


def extract_office_text(path: Path, extension: str) -> tuple[str, str]:
    normalized = extension.lower()
    if normalized == ".docx":
        text, _ = extract_docx_text(path)
        return text, "docx"
    if normalized == ".pptx":
        text, _ = extract_pptx_text(path)
        return text, "pptx"
    if normalized in LEGACY_OFFICE_EXTENSIONS:
        raise HTTPException(
            status_code=415,
            detail=f"{normalized} is a legacy binary Office format. Upload {normalized}x or convert it to PDF first.",
        )
    raise HTTPException(status_code=415, detail=f"Unsupported material format: {normalized or 'unknown'}")
